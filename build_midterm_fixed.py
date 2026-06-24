#!/usr/bin/env python3
"""
Build midterm PPT from official template:
- Preserve backgrounds 100%
- Short bullets (no overflow)
- Times New Roman 12pt body, 24pt titles
- Embedded speaker notes
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
import shutil
import os
import zipfile

TEMPLATE = os.path.expanduser("~/Desktop/副本Generation-AI 中期汇报模版.pptx")
PROJECT = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(PROJECT, "midterm-final.pptx")
FIXED_OUT = os.path.expanduser("~/Desktop/未来项目/Reframe-Destiny-Midterm-FIXED.pptx")
DESKTOP_OUT = os.path.expanduser("~/Desktop/Reframe-Destiny-Midterm-FIXED.pptx")

NAME = "Xinyun Zhang (张馨云)"
ADVISOR = "Lawted Wu"
TITLE = "Reframe Destiny: Gender Bias in BaZi & Astrology"

FONT = "Times New Roman"
BODY_PT = 12
TITLE_PT = 24

# Short body text — max ~5 bullets, ~12 words each
BODIES = {
    3: (
        "Research Question:\n"
        "How does an AI website help young women find and reframe gender bias in BaZi and astrology?\n\n"
        "Hypothesis:\n"
        "After using Reframe Destiny, users will notice bias more and feel more confident to rewrite their stories."
    ),
    4: (
        "• Young women meet BaZi and astrology through family and social media.\n"
        "• Same chart: \"leader\" for men, \"too strong\" for women.\n"
        "• We study narrative bias — not whether divination is true.\n"
        "• Question: Who narrates your fate?\n"
        "• Reframe Destiny is a critical tool, not a fortune-telling app."
    ),
    5: (
        "• Butler (1990): Gender is built through cultural stories.\n"
        "• Foucault (1972): Language has power to shape people.\n"
        "• Haraway (1988): Who reads the chart changes the meaning.\n"
        "• Bender et al. (2021): AI can help critique, not predict fate.\n"
        "• Gap: No tool compares gender views on the same chart."
    ),
    6: (
        "• Artifact: bilingual website, 7-step Journey + Bias Scanner.\n"
        "• Plan: 15–25 young women; pre/post survey; 5 interviews.\n"
        "• Code 40 traditional readings for 7 bias types.\n"
        "• User study NOT run yet — design only at midterm."
    ),
    7: (
        "• Done: idea, research folder, interactive prototype.\n"
        "• Next: user study (late June); final paper (July).\n"
        "• Challenge: missed 4–5 classes; catching up now.\n"
        "• Plan: add gender-switch level after midterm."
    ),
    8: (
        "Expected Results — user study not yet run.\n\n"
        "• Users may notice marriage bias and gender stereotypes most.\n"
        "• Bias-awareness scores may increase after one session.\n"
        "• Reframed reading may feel more empowering than traditional text.\n"
        "• [Insert prototype screenshot here]"
    ),
    9: (
        "Bender, E. M., et al. (2021). Stochastic parrots. FAccT.\n"
        "Butler, J. (2006). Gender trouble. Routledge.\n"
        "Foucault, M. (1972). Archaeology of knowledge. Pantheon.\n"
        "Haraway, D. (1988). Situated knowledges. Feminist Studies.\n"
        "Vyse, S. (2018). Believing in magic. Oxford."
    ),
}

TITLE_OVERRIDES = {
    8: "Preliminary Results\n(study not run yet)",
}

# Speaker notes — Grade 10 English, ~7 min total
NOTES = [
    # Slide 1 Title
    "Hi everyone. My name is Xinyun Zhang.\n\n"
    "Today I will present Reframe Destiny — a website that helps young women find gender bias in BaZi and astrology, and rewrite those stories.\n\n"
    "My advisor is Lawted Wu. This is my Generation AI Social Science project.",

    # Slide 2 RQ
    "My research question is: How does an AI website help young women find and reframe gender bias in BaZi and astrology readings?\n\n"
    "My hypothesis is: After using Reframe Destiny, users will notice bias more easily and feel more confident to rewrite their stories.",

    # Slide 3 Background
    "Many young women learn BaZi or astrology from family and social media.\n\n"
    "These systems tell stories about marriage and life paths — not just personality.\n\n"
    "I study BaZi myself. The same chart can mean 'leader' for men but 'too strong' for women.\n\n"
    "We do not ask if divination is true. We ask: who tells your story, and is it unfair to women?\n\n"
    "Reframe Destiny is not a fortune-telling app. It is a tool to think critically.",

    # Slide 4 Literature
    "Butler says gender is built through culture and stories.\n\n"
    "Foucault says language has power.\n\n"
    "Haraway says who reads the data matters.\n\n"
    "Bender warns AI can copy bias — here AI helps us question, not predict fate.\n\n"
    "The gap: no website lets users compare gender views on the same chart. That is what I built.",

    # Slide 5 Method
    "My project is a bilingual website with seven steps.\n\n"
    "Users pick BaZi or astrology, read a traditional version, use a Bias Scanner, and see an AI Reframe.\n\n"
    "I plan to test 15 to 25 young women with surveys and short interviews.\n\n"
    "Important: I have NOT done the user study yet. Today I show my design only.",

    # Slide 6 Plan
    "Timeline: idea and prototype — done. User study — late June. Final paper — July.\n\n"
    "My challenge: I missed several classes and I am catching up.\n\n"
    "After midterm I will run the study and add a gender-switch feature.",

    # Slide 7 Results
    "These are expected results. I have NOT run the study yet.\n\n"
    "I think users will notice marriage bias and gender stereotypes most often.\n\n"
    "I think the reframed reading will feel more empowering.\n\n"
    "Point to your screenshot here.",

    # Slide 8 References
    "My references are Butler, Foucault, Haraway, Bender, and Vyse — listed on this slide in APA format.",

    # Slide 9 Thank you
    "Thank you for your feedback. I welcome your questions.",
]


def apply_font(paragraph, size_pt, bold=False):
    for run in paragraph.runs:
        run.font.name = FONT
        run.font.size = Pt(size_pt)
        run.font.bold = bold


def set_multiline_text(shape, text, size_pt=BODY_PT):
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    lines = text.split("\n")
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, line in enumerate(lines):
        p = tf.paragraphs[i]
        p.text = line
        p.level = 0
        p.line_spacing = 1.15
        apply_font(p, size_pt)
    for i in range(len(lines), len(tf.paragraphs)):
        tf.paragraphs[i].text = ""


def set_single_text(shape, text, size_pt=TITLE_PT, bold=True):
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    if tf.paragraphs:
        lines = text.split("\n")
        tf.paragraphs[0].text = lines[0]
        apply_font(tf.paragraphs[0], size_pt, bold=bold)
        for j, line in enumerate(lines[1:], 1):
            if j < len(tf.paragraphs):
                tf.paragraphs[j].text = line
                apply_font(tf.paragraphs[j], size_pt, bold=bold)
            else:
                p = tf.add_paragraph()
                p.text = line
                apply_font(p, size_pt, bold=bold)
        for p in tf.paragraphs[len(lines):]:
            p.text = ""


def remove_first_slide(prs):
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]


def fill_slide(slide, orig_num):
    shapes = slide.shapes
    if orig_num == 2:
        if len(shapes) > 2 and shapes[2].has_text_frame:
            set_single_text(shapes[2], TITLE, size_pt=22)
        if len(shapes) > 3 and shapes[3].has_text_frame:
            set_multiline_text(
                shapes[3],
                f"Student Name: {NAME}\nResearch Advisor: {ADVISOR}\nGeneration AI 2026 · Social Science Track",
                size_pt=BODY_PT,
            )
        return
    if orig_num == 10:
        if len(shapes) > 1 and shapes[1].has_text_frame:
            set_single_text(shapes[1], "", size_pt=BODY_PT, bold=False)
        if len(shapes) > 4 and shapes[4].has_text_frame:
            set_multiline_text(
                shapes[4],
                f"Student Name: {NAME}\nResearch Advisor: {ADVISOR}",
                size_pt=BODY_PT,
            )
        return
    if orig_num in TITLE_OVERRIDES and len(shapes) > 0 and shapes[0].has_text_frame:
        set_single_text(shapes[0], TITLE_OVERRIDES[orig_num], size_pt=20)
    if orig_num in BODIES and len(shapes) > 3 and shapes[3].has_text_frame:
        set_multiline_text(shapes[3], BODIES[orig_num], size_pt=BODY_PT)


def add_notes(slide, note_text):
    if not note_text:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = note_text


def main():
    if not os.path.exists(TEMPLATE):
        print(f"ERROR: Template not found: {TEMPLATE}")
        return 1

    os.makedirs(os.path.dirname(FIXED_OUT), exist_ok=True)

    shutil.copy2(TEMPLATE, OUTPUT)
    prs = Presentation(OUTPUT)
    remove_first_slide(prs)

    for idx, orig_num in enumerate(range(2, 11)):
        if idx < len(prs.slides):
            fill_slide(prs.slides[idx], orig_num)
            if idx < len(NOTES):
                add_notes(prs.slides[idx], NOTES[idx])

    prs.save(OUTPUT)
    shutil.copy2(OUTPUT, FIXED_OUT)
    shutil.copy2(OUTPUT, DESKTOP_OUT)

    with zipfile.ZipFile(OUTPUT) as z:
        media = len([n for n in z.namelist() if n.startswith("ppt/media/") and not n.endswith("/")])
    print(f"Created: {FIXED_OUT}")
    print(f"Also: {DESKTOP_OUT}")
    print(f"Slides: {len(prs.slides)}, Media: {media}, Notes embedded: {len(NOTES)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
