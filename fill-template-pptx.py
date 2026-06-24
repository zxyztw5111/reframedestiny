#!/usr/bin/env python3
"""
Fill official Generation AI midterm template — TEXT ONLY.
Preserves all backgrounds, images, layouts, and formatting.
Only replaces text inside existing text boxes.
"""
from pptx import Presentation
import shutil
import os
import zipfile

TEMPLATE = os.path.expanduser("~/Desktop/副本Generation-AI 中期汇报模版.pptx")
OUTPUT = os.path.join(os.path.dirname(__file__), "midterm-final.pptx")
DESKTOP_OUT = os.path.expanduser("~/Desktop/Reframe-Destiny-Midterm.pptx")

NAME = "Xinyun Zhang (张馨云)"
ADVISOR = "Lawted Wu"
TITLE = (
    "Reframe Destiny: Detecting and Reframing Gendered "
    "Narrative Bias in BaZi and Western Astrology"
)

# Body text for content box (shape index 3 on slides 3–9 in original template)
BODIES = {
    3: (  # Research Question
        "Research Question:\n"
        "How does an AI-assisted interactive website affect young women's ability "
        "to identify and reframe gendered narrative bias in traditional BaZi and "
        "Western astrological readings?\n\n"
        "Hypothesis:\n"
        "I hypothesize that after using Reframe Destiny, young women will show "
        "increased awareness of gender bias in divination narratives and greater "
        "confidence in reframing these narratives for themselves."
    ),
    4: (  # Background
        "• Many young women encounter BaZi or astrology through family, social media, or apps.\n"
        "• These systems tell stories about marriage, morality, and life paths.\n"
        "• The same chart element can mean \"ambitious leader\" for men but "
        "\"too strong\" or \"harmful to husband\" for women.\n"
        "• This project does NOT ask whether divination is true.\n"
        "• It asks: who narrates your fate, and does that narration carry gender bias?\n"
        "• Reframe Destiny is a critical interactive tool — not a fortune-telling app."
    ),
    5: (  # Literature Review
        "• Butler (1990): Gender is performed through cultural narratives.\n"
        "• Foucault (1972): Discourse shapes what counts as knowledge.\n"
        "• Haraway (1988): Who interprets data matters — same chart, different meanings.\n"
        "• Bender et al. (2021): AI reproduces bias; here AI is a lens for critique.\n"
        "• Gap: No tool lets users compare gender perspectives and reframe narratives.\n"
        "• My contribution: cross-cultural (BaZi + astrology) + interactive bias discovery."
    ),
    6: (  # Research Design / Method
        "Artifact: Reframe Destiny — bilingual website, 7-step Journey.\n"
        "Planned participants: 15–25 young women, age 17–34.\n"
        "Procedure: Pre-survey → 15-min session → post-survey → 5 interviews.\n"
        "Analysis: Pre/post Likert scores; bias category frequency.\n"
        "Content analysis: Coding 40 traditional readings for 7 bias types.\n"
        "Note: User study NOT yet run — midterm shows design only."
    ),
    7: (  # Research Plan and Challenges
        "Timeline:\n"
        "• Week 2: Idea + GitHub — Done\n"
        "• Week 3–4: Research archive + prototype — Done\n"
        "• Midterm (June 25): Defense — In progress\n"
        "• Late June: User study (n=15–25) — Planned\n"
        "• July: Final paper (3000–5000 words, APA) — Planned\n\n"
        "Challenges: Missed 4–5 classes; user study not started.\n"
        "Plan: Run user study after midterm; add gender-switch level to Journey."
    ),
    8: (  # Preliminary Results
        "Expected Results — user study not yet run\n\n"
        "I have NOT collected user data yet. Expected outcomes:\n"
        "• Users will most often identify marriage centrism and gender stereotypes.\n"
        "• Pre/post bias-awareness scores will increase after one session.\n"
        "• Side-by-side traditional vs. reframed reading will feel more empowering.\n"
        "• Informal testing (n=1–2): 7-step flow works; 3 bias fragments unlock per journey.\n\n"
        "[Insert prototype screenshot here]"
    ),
    9: (  # References
        "Bender, E. M., et al. (2021). On the dangers of stochastic parrots. FAccT, 610–623.\n\n"
        "Butler, J. (2006). Gender trouble. Routledge.\n\n"
        "Foucault, M. (1972). The archaeology of knowledge. Pantheon.\n\n"
        "Haraway, D. (1988). Situated knowledges. Feminist Studies, 14(3), 575–599.\n\n"
        "Vyse, S. (2018). Believing in magic. Oxford University Press."
    ),
}

TITLE_OVERRIDES = {
    8: "Preliminary Results — Expected Results — user study not yet run",
}


def set_multiline_text(shape, text):
    """Replace text in existing text frame without clearing shape or layout."""
    tf = shape.text_frame
    lines = text.split("\n")
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, line in enumerate(lines):
        tf.paragraphs[i].text = line
    for i in range(len(lines), len(tf.paragraphs)):
        tf.paragraphs[i].text = ""


def set_single_text(shape, text):
    """Replace single-line text preserving text box."""
    tf = shape.text_frame
    if tf.paragraphs:
        tf.paragraphs[0].text = text
        for p in tf.paragraphs[1:]:
            p.text = ""


def remove_first_slide(prs):
    """Remove slide 1 (instructions page) per teacher requirement."""
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]


def fill_slide(slide, orig_num):
    """Fill one slide based on original template slide number (2–10)."""
    shapes = slide.shapes

    if orig_num == 2:
        # Title page: shape[2]=title, shape[3]=name/advisor
        if len(shapes) > 2 and shapes[2].has_text_frame:
            set_single_text(shapes[2], TITLE)
        if len(shapes) > 3 and shapes[3].has_text_frame:
            set_multiline_text(
                shapes[3],
                f"Student Name: {NAME}\nResearch Advisor: {ADVISOR}\nGeneration AI 2026 · Social Science Track",
            )
        return

    if orig_num == 10:
        # Thank you page
        if len(shapes) > 1 and shapes[1].has_text_frame:
            set_single_text(shapes[1], "")  # clear Chinese line
        if len(shapes) > 4 and shapes[4].has_text_frame:
            set_multiline_text(
                shapes[4],
                f"Student Name: {NAME}\nResearch Advisor: {ADVISOR}",
            )
        return

    # Slides 3–9: shape[0]=title, shape[3]=content box
    if orig_num in TITLE_OVERRIDES:
        if len(shapes) > 0 and shapes[0].has_text_frame:
            set_single_text(shapes[0], TITLE_OVERRIDES[orig_num])

    if orig_num in BODIES:
        if len(shapes) > 3 and shapes[3].has_text_frame:
            set_multiline_text(shapes[3], BODIES[orig_num])


def count_media(pptx_path):
    with zipfile.ZipFile(pptx_path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/") and not n.endswith("/")]
        slides = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
    return len(media), len(slides)


def verify_pptx(template_path, output_path):
    orig_media, orig_slides = count_media(template_path)
    out_media, out_slides = count_media(output_path)
    # After removing 1 slide: slides should be orig_slides - 1
    ok_media = out_media == orig_media
    ok_slides = out_slides == orig_slides - 1
    return {
        "template_media": orig_media,
        "output_media": out_media,
        "template_slides": orig_slides,
        "output_slides": out_slides,
        "media_preserved": ok_media,
        "slide_count_ok": ok_slides,
    }


def main():
    if not os.path.exists(TEMPLATE):
        print(f"ERROR: Template not found: {TEMPLATE}")
        return 1

    shutil.copy2(TEMPLATE, OUTPUT)
    prs = Presentation(OUTPUT)

    remove_first_slide(prs)

    # Original slides 2–10 → after removal, index 0–8 maps to orig 2–10
    for idx, orig_num in enumerate(range(2, 11)):
        if idx < len(prs.slides):
            fill_slide(prs.slides[idx], orig_num)

    prs.save(OUTPUT)

    # Copy to Desktop
    shutil.copy2(OUTPUT, DESKTOP_OUT)

    report = verify_pptx(TEMPLATE, OUTPUT)
    print(f"Created: {OUTPUT}")
    print(f"Desktop: {DESKTOP_OUT}")
    print(f"Slides: {report['output_slides']} (expected {report['template_slides'] - 1})")
    print(f"Media files: {report['output_media']} (template had {report['template_media']})")
    print(f"Media preserved: {report['media_preserved']}")
    print(f"Slide count OK: {report['slide_count_ok']}")

    if not report["media_preserved"]:
        print("WARNING: Media count mismatch — backgrounds may be affected!")
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
